"""Generate the Resume Screener MCP talk deck as an editable PowerPoint file.

python-pptx is NOT a project dependency. Run this with an ephemeral environment:

    uv run --with python-pptx python talk/deck/build_deck.py

Output: talk/deck/presentation.pptx (~16 slides, each with speaker notes).
The deck uses only built-in shapes/fonts (no external assets), so it builds anywhere.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Inches

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
INK = RGBColor(0x1A, 0x1A, 0x2E)        # near-black text
ACCENT = RGBColor(0x2D, 0x6C, 0xDF)     # blue accent
ACCENT_DARK = RGBColor(0x16, 0x2B, 0x57)  # deep blue (divider bg)
MUTED = RGBColor(0x5A, 0x5A, 0x6E)      # muted grey for sub-text
BG = RGBColor(0xF6, 0xF7, 0xFB)         # light page background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD9, 0xDD, 0xE8)       # hairline rule

FONT = "Calibri"

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank(prs: Presentation):
    """Add a blank slide and paint the page background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back so everything else paints on top
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return slide


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(run, *, size, bold=False, color=INK, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def _notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def _accent_bar(slide, top=Inches(1.5), left=Inches(0.9), width=Inches(0.55)):
    bar = slide.shapes.add_shape(1, left, top, width, Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def title_slide(prs, *, title, subtitle, presenter, notes=""):
    slide = _blank(prs)
    # dark band
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_DARK
    band.line.fill.background()
    band.shadow.inherit = False

    _, tf = _textbox(slide, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, size=46, bold=True, color=WHITE)

    p2 = tf.add_paragraph(); p2.space_before = Pt(14)
    r2 = p2.add_run(); r2.text = subtitle
    _set(r2, size=22, color=RGBColor(0xBF, 0xD0, 0xF2))

    _, tf3 = _textbox(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.8))
    r3 = tf3.paragraphs[0].add_run(); r3.text = presenter
    _set(r3, size=16, color=RGBColor(0x9A, 0xB0, 0xDC))

    # accent underline
    bar = slide.shapes.add_shape(1, Inches(1.0), Inches(2.25), Inches(2.2), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    _notes(slide, notes)
    return slide


def content_slide(prs, *, number, title, bullets, notes=""):
    """bullets: list of (text, level) or (text, level, kind) where kind in
    {'', 'key', 'code'}. level 0 = top bullet, 1 = sub-bullet."""
    slide = _blank(prs)

    # kicker (slide number / section)
    _, tfk = _textbox(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5))
    rk = tfk.paragraphs[0].add_run(); rk.text = f"{number:02d}"
    _set(rk, size=14, bold=True, color=ACCENT)

    # title
    _, tft = _textbox(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(1.0))
    rt = tft.paragraphs[0].add_run(); rt.text = title
    _set(rt, size=32, bold=True, color=INK)

    _accent_bar(slide, top=Inches(1.85))

    # body
    _, tfb = _textbox(slide, Inches(0.95), Inches(2.2), Inches(11.4), Inches(4.9))
    first = True
    for item in bullets:
        text, level = item[0], item[1]
        kind = item[2] if len(item) > 2 else ""
        p = tfb.paragraphs[0] if first else tfb.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8 if level == 0 else 4)
        # bullet marker
        marker = "" if kind in ("key",) else ("–  " if level else "•  ")
        if kind == "key":
            marker = "★  "
        r = p.add_run(); r.text = marker + text
        if kind == "key":
            _set(r, size=19, bold=True, color=ACCENT)
        elif kind == "code":
            _set(r, size=16, color=INK, font="Consolas")
        elif level == 0:
            _set(r, size=20, color=INK)
        else:
            _set(r, size=17, color=MUTED)
    _notes(slide, notes)
    return slide


def two_col_slide(prs, *, number, title, left_head, left_items, right_head,
                  right_items, notes="", left_accent=MUTED, right_accent=ACCENT):
    slide = _blank(prs)
    _, tfk = _textbox(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5))
    rk = tfk.paragraphs[0].add_run(); rk.text = f"{number:02d}"
    _set(rk, size=14, bold=True, color=ACCENT)
    _, tft = _textbox(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(1.0))
    rt = tft.paragraphs[0].add_run(); rt.text = title
    _set(rt, size=32, bold=True, color=INK)
    _accent_bar(slide, top=Inches(1.85))

    def col(left, head, items, head_color):
        card = slide.shapes.add_shape(1, left, Inches(2.25), Inches(5.55), Inches(4.6))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RULE; card.line.width = Pt(1)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.35); tf.margin_right = Inches(0.35)
        tf.margin_top = Inches(0.3)
        hp = tf.paragraphs[0]
        hr = hp.add_run(); hr.text = head
        _set(hr, size=20, bold=True, color=head_color)
        hp.space_after = Pt(10)
        for text, *rest in items:
            lvl = rest[0] if rest else 0
            p = tf.add_paragraph(); p.level = lvl
            p.space_after = Pt(6)
            r = p.add_run(); r.text = ("–  " if lvl else "•  ") + text
            _set(r, size=16 if lvl == 0 else 14, color=INK if lvl == 0 else MUTED)

    col(Inches(0.95), left_head, left_items, left_accent)
    col(Inches(6.85), right_head, right_items, right_accent)
    _notes(slide, notes)
    return slide


def divider_slide(prs, *, title, subtitle, notes=""):
    slide = _blank(prs)
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT_DARK
    band.line.fill.background(); band.shadow.inherit = False
    tb, tf = _textbox(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(2.0))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    _set(r, size=54, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.space_before = Pt(12)
    r2 = p2.add_run(); r2.text = subtitle
    _set(r2, size=22, color=RGBColor(0xBF, 0xD0, 0xF2))
    bar = slide.shapes.add_shape(1, Inches(1.0), Inches(2.75), Inches(2.2), Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    _notes(slide, notes)
    return slide


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Title
    title_slide(
        prs,
        title="Resume Screener MCP",
        subtitle="Screening resumes with Claude — honestly.",
        presenter="Omkar Jadhav  ·  NonStop.io",
        notes=(
            "Intro (~20s). I want to show you a small but complete tool — a Resume Screener "
            "MCP server. It ranks a folder of resumes against a job description, with reasons, "
            "and does it honestly. By the end I'll send a real interview email, live."
        ),
    )

    # 2 — The problem
    content_slide(
        prs, number=2, title="The problem",
        bullets=[
            ("Recruiters screen hundreds of resumes by hand.", 0),
            ("Slow — hours per role.", 1),
            ("Inconsistent — the 200th resume gets a different read than the 1st.", 1),
            ("Lossy — scanned PDFs, locked files, legacy .doc silently fail to open.", 1),
            ("A dropped file is a missed candidate.", 0),
            ("Silent data loss is the cardinal sin of a screening tool.", 0, "key"),
        ],
        notes=(
            "~60s. Three problems: slow, inconsistent, and lossy. Unreadable files silently "
            "vanish from the pile. Land the line: silent data loss is the cardinal sin."
        ),
    )

    # 3 — What is an MCP server
    content_slide(
        prs, number=3, title="What is an MCP server?",
        bullets=[
            ("MCP = Model Context Protocol: it gives Claude tools it can call.", 0),
            ("The server provides capabilities + data.", 0),
            ("The model provides the intelligence.", 0),
            ("Runs as a small subprocess; Claude Code talks to it over stdio (JSON-RPC).", 1),
            ("I ask Claude to 'screen these resumes' — it decides which tools to call.", 1),
        ],
        notes=(
            "~90s. Context for newer folks. The split that matters: server = capabilities + "
            "data, model = intelligence. Skip/trim this slide if the room knows MCP."
        ),
    )

    # 4 — Solution at a glance
    content_slide(
        prs, number=4, title="The solution at a glance",
        bullets=[
            ("Point Claude at a folder + a job description.", 0),
            ("Get back a ranked shortlist — scores + human-readable reasons.", 0),
            ("Then optional actions:", 0),
            ("export to CSV/Excel  ·  collect chosen resumes  ·  draft + send emails.", 1),
            ("Pipeline:  extract  →  pre-filter (if needed)  →  Claude judges  →  act.", 0, "code"),
        ],
        notes=(
            "~60s. Point at a folder + JD, get a ranked shortlist with reasons, then optional "
            "actions. Walk the pipeline line left to right."
        ),
    )

    # 5 — Claude is the judge (two-column A vs B)
    two_col_slide(
        prs, number=5, title='Key idea: "Claude is the judge"',
        left_head="Model A — server calls an LLM",
        left_items=[
            ("Server hits an LLM API, returns scores.", 0),
            ("Needs an API key.", 0),
            ("Per-run cost + rate limits.", 0),
            ("Redundant — you're already in Claude.", 0),
        ],
        right_head="Model B — what I built ✓",
        right_items=[
            ("Server prepares a 'judging packet'.", 0),
            ("JD + rubric + cleaned resume text.", 1),
            ("Claude scores 0–100 in-session.", 0),
            ("No API key. No cost. No rate limit.", 0),
            ("Idiomatic MCP.", 0),
        ],
        notes=(
            "~90s. The most important decision. Model A: server calls an API — key, cost, "
            "limits, redundant. Model B: server preps data, Claude judges in-session. "
            "Land the line: the judge is the model you already have."
        ),
    )

    # 6 — Two-stage pipeline
    content_slide(
        prs, number=6, title="A two-stage, deterministic pipeline",
        bullets=[
            ("Stage 1 — Extraction (pure Python).", 0),
            ("pypdf + python-docx; normalize text; quality gate ≥ 100 chars.", 1),
            ("The gate catches scanned PDFs with no real text.", 1),
            ("Stage 2 — Scale.", 0),
            ("≤ 25 readable → hand them all to Claude.", 1),
            ("> 25 → local TF-IDF pre-filter narrows to ~25 first.", 1),
            ("Coarse on purpose — and we tell you how many were set aside.", 1),
        ],
        notes=(
            "~75s. Stage 1 extraction with a 100-char quality gate. Stage 2 scale: small pile "
            "goes straight to Claude; big pile gets a cheap TF-IDF pre-filter, and we surface "
            "how many we set aside."
        ),
    )

    # 7 — Tools: screening
    two_col_slide(
        prs, number=7, title="The 11 tools — screening",
        left_head="Inspect",
        left_items=[
            ("ping — health check.", 0),
            ("list_resumes — inventory a folder; show failures.", 0),
        ],
        right_head="Rank",
        right_items=[
            ("screen_resumes — small piles.", 0),
            ("screen_resumes_bulk — large piles (pre-filter).", 0),
            ("compare_candidates — 2–4 finalists side by side.", 0),
            ("rerank — re-score under a new priority.", 0),
        ],
        notes=(
            "~45s. Six screening tools. rerank lets you say 'now weight AWS more' without "
            "re-pasting the JD."
        ),
    )

    # 8 — Tools: actions
    two_col_slide(
        prs, number=8, title="The 11 tools — actions",
        left_head="Output",
        left_items=[
            ("export_shortlist — write CSV / Excel.", 0),
            ("collect_selected — copy chosen resumes to a folder.", 0),
        ],
        right_head="Email",
        right_items=[
            ("get_email_templates — starter templates.", 0),
            ("draft_emails — editable .eml drafts on disk.", 0),
            ("send_emails — preview, then send (gated).", 0),
        ],
        notes=(
            "~45s. Five action tools you take AFTER the ranking. We'll see the email three "
            "carefully in a minute."
        ),
    )

    # 9 — Fairness by construction
    content_slide(
        prs, number=9, title="Fairness by construction",
        bullets=[
            ("Anonymization — redact BEFORE the model sees the text.", 0),
            ("Telling an LLM 'ignore the name' doesn't work — it already saw it.", 1),
            ("Server strips name/email/phone/school; labels people candidate_01.", 1),
            ("Knockouts — hard rules applied by code, not by the model.", 0),
            ("'5+ years', 'must list Python' → deterministic, auditable verdicts.", 1),
            ("Honest: reduces direct-signal bias; doesn't remove writing-style signals.", 0, "key"),
        ],
        notes=(
            "~90s. Two structural fairness features. Anonymization must happen in the server, "
            "before text reaches the model. Knockouts are deterministic code, not fuzzy "
            "instructions. Be explicit about the limit: reduces, not eliminates."
        ),
    )

    # 10 — Never silently drop a file
    content_slide(
        prs, number=10, title="Never silently drop a file",
        bullets=[
            ("Every unreadable or knocked-out file appears in a list — with a reason.", 0),
            ('"Found 142, parsed 138, 4 failed — and here is why each failed."', 0, "code"),
            ("You always know the true denominator.", 0),
            ("Nothing disappears quietly.", 0, "key"),
        ],
        notes=(
            "~45s. Back to the cardinal sin. You always know the true denominator. Nothing "
            "disappears quietly."
        ),
    )

    # 11 — The email workflow
    content_slide(
        prs, number=11, title="The email workflow",
        bullets=[
            ("Pull a template  →  personalize per candidate.", 0),
            ("Write each as an editable .eml file on disk.", 0),
            ("A human can open and edit it; the edit wins at send time.", 1),
            ("Preview (dry-run)  →  send.", 0),
        ],
        notes=(
            "~45s. Templates → personalized drafts as real files on disk → preview → send. "
            "Drafts are editable; hand-edits win."
        ),
    )

    # 12 — Email safety model
    content_slide(
        prs, number=12, title="Email safety model",
        bullets=[
            ("send_emails is dry-run by default — shows what WOULD go out, sends nothing.", 0),
            ("Real send needs TWO flags:  dry_run=False  AND  confirm=True.", 0, "code"),
            ("Credentials come from env vars — never logged.", 0),
            ("Every real send is appended to send_log.json (audit trail).", 0),
            ("Nothing irreversible happens by accident.", 0, "key"),
        ],
        notes=(
            "~60s. Email is the one irreversible, outward-facing action — so it's gated. "
            "Dry-run default; real send needs both flags. Creds never logged; sends logged. "
            "Land the line."
        ),
    )

    # 13 — Tech stack & how it runs
    two_col_slide(
        prs, number=13, title="Tech stack & how it runs",
        left_head="Stack",
        left_items=[
            ("Python 3.11+, packaged with uv.", 0),
            ("FastMCP over stdio.", 0),
            ("scikit-learn (TF-IDF).", 0),
            ("pypdf, python-docx, openpyxl.", 0),
        ],
        right_head="Why it's solid",
        right_items=[
            ("One-command install via uvx.", 0),
            ("102 tests, all deterministic.", 0),
            ("No LLM in the server → fully unit-testable.", 0),
            ("That's the payoff of Model B.", 0),
        ],
        notes=(
            "~45s. For the engineers. The big one: no LLM in the server means the whole thing "
            "is deterministic and unit-testable — 102 tests. Then: enough slides, let me show you."
        ),
    )

    # 14 — Live demo divider
    divider_slide(
        prs, title="Live demo",
        subtitle="From a folder of resumes → a real email. (Runbook: demo/DEMO_GUIDE.md)",
        notes=(
            "~5–6 min. Switch to Claude Code; follow DEMO_GUIDE.md. "
            "1) list_resumes — show failures with reasons. "
            "2) screen — server returned no scores; Claude judged. "
            "3) (optional) compare/rerank. "
            "4) collect top 3. "
            "5) draft invites — all to my own inbox; open one .eml. "
            "6) dry-run preview — the safety gate. "
            "7) real send — switch to inbox, show it arrived; mention send_log.json. "
            "Land: screened, ranked, drafted, and sent without leaving the chat."
        ),
    )

    # 15 — Honest limitations
    content_slide(
        prs, number=15, title="Honest limitations",
        bullets=[
            ("TF-IDF pre-filter is coarse — a borderline resume can be set aside.", 0),
            ("We surface the count so you can review them.", 1),
            ("Anonymization handles direct identifiers, not every indirect signal.", 0),
            ("Knockouts use heuristics — when ambiguous, we warn instead of dropping.", 0),
            ("OCR for scanned PDFs is built but optional and off by default.", 0),
            ("None of this is hidden — it's in the tool outputs and the docs.", 0, "key"),
        ],
        notes=(
            "~60s. Honesty is the whole theme, so be straight about trade-offs. Every one of "
            "these is surfaced in outputs/docs, not buried."
        ),
    )

    # 16 — Takeaways + Q&A
    content_slide(
        prs, number=16, title="Takeaways",
        bullets=[
            ("Transparent — every failure surfaced.", 0),
            ("Reproducible — deterministic parts give the same answer every time.", 0),
            ("Auditable — knockouts have reasons; sends have logs.", 0),
            ("The human stays in control the whole way.", 0),
            ("The judge is the model you were already talking to.", 0, "key"),
            ("Thank you — questions?", 0),
        ],
        notes=(
            "~60s. Four takeaways + the elegant punchline. Then Q&A. Likely Qs: accuracy "
            "(scales with the model; deterministic parts reproducible), why not API in server "
            "(cost/keys/limits/redundancy), bias (reduces not cures), scanned resumes (optional OCR)."
        ),
    )

    out = Path(__file__).with_name("presentation.pptx")
    prs.save(out)
    return out, len(prs.slides)


if __name__ == "__main__":
    path, n = build()
    print(f"Wrote {path}  ({n} slides)")
